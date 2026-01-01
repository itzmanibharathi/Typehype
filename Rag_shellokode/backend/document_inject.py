import os
import logging
import requests  # ← Added for OCR.space API
from pptx import Presentation
from docx import Document
import pandas as pd
import pdfplumber
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ==================== TEXT EXTRACTION FUNCTIONS ====================

def extract_from_pdf(file_path: str) -> str:
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
    return text.strip()


def extract_from_pptx(file_path: str) -> str:
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
    return text.strip()


def extract_from_docx(file_path: str) -> str:
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
    return text.strip()


def extract_from_excel(file_path: str) -> str:
    text = ""
    try:
        sheets = pd.read_excel(file_path, sheet_name=None)
        for name, df in sheets.items():
            text += f"Sheet: {name}\n{df.to_string(index=False)}\n\n"
    except Exception as e:
        logger.error(f"Error extracting Excel: {e}")
    return text.strip()


def extract_from_txt(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()
    except Exception as e:
        logger.error(f"Error extracting TXT: {e}")
        return ""


def extract_from_csv(file_path: str) -> str:
    try:
        df = pd.read_csv(file_path)
        return df.to_string(index=False).strip()
    except Exception as e:
        logger.error(f"Error extracting CSV: {e}")
        return ""


def extract_from_image(file_path: str) -> str:
    """Extract text from image using OCR.space API (free, no install needed)"""
    try:
        url = "https://api.ocr.space/parse/image"
        
        with open(file_path, "rb") as image_file:
            files = {"file": image_file}
            data = {
                "apikey": "helloworld",          # Free default key for testing
                "language": "eng",               # Change to 'chi', 'fra', etc. if needed
                "isOverlayRequired": False,
                # "filetype": "JPG"              # Optional – usually auto-detected
            }
            
            response = requests.post(url, files=files, data=data, timeout=30)
            response.raise_for_status()
            result = response.json()
            
            if result.get("ParsedResults"):
                text = result["ParsedResults"][0].get("ParsedText", "")
                return text.strip()
            else:
                error = result.get("ErrorMessage", ["Unknown error"])
                logger.warning(f"OCR.space error for {file_path}: {error}")
                return ""
                
    except requests.exceptions.RequestException as e:
        logger.warning(f"Network/error extracting text from image {file_path}: {e}")
        return ""
    except Exception as e:
        logger.warning(f"Failed to extract text from image {file_path}: {e}")
        return ""


def extract_text(file_path: str, filename: str = "") -> str:
    """Main function to extract text from any supported file"""
    ext = os.path.splitext(filename or file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_from_pptx(file_path)
    elif ext == ".docx":
        return extract_from_docx(file_path)
    elif ext in {".xlsx", ".xls"}:
        return extract_from_excel(file_path)
    elif ext == ".txt":
        return extract_from_txt(file_path)
    elif ext == ".csv":
        return extract_from_csv(file_path)
    
    # Image files – using OCR.space (no pytesseract!)
    elif ext in {".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff", ".webp"}:
        return extract_from_image(file_path)  # Always returns a string
    
    else:
        logger.warning(f"Unsupported file type: {ext}")
        return ""


def chunk_text(text: str) -> list[str]:
    """Split text into chunks for embedding"""
    if not text.strip():
        return []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
        chunk_overlap=200,
        length_function=len,
    )
    return splitter.split_text(text)
